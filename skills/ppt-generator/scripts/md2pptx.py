import sys
import re
import os
import yaml
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip('#')
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))

def apply_text_style(shape, style_config, run_type='body'):
    if not shape.has_text_frame:
        return
        
    tf = shape.text_frame
    
    # Get configuration for this run type
    if run_type == 'title':
        font_config = style_config.get('typography', {}).get('heading', {})
        color_hex = style_config.get('colors', {}).get('text', '#000000') # Default title logic usually matches body or accent
        # Allow specific override if exists, else match general text
    else:
        font_config = style_config.get('typography', {}).get('body', {})
        color_hex = style_config.get('colors', {}).get('text', '#000000')

    font_name = font_config.get('font', None)
    
    # Iterate paragraphs and runs
    for p in tf.paragraphs:
        for run in p.runs:
            if font_name:
                run.font.name = font_name
            
            if color_hex:
                 try:
                     run.font.color.rgb = hex_to_rgb(color_hex)
                 except:
                     pass

def apply_slide_background(slide, style_config):
    bg_color = style_config.get('colors', {}).get('background', None)
    if bg_color:
        background = slide.background
        fill = background.fill
        fill.solid()
        try:
            fill.fore_color.rgb = hex_to_rgb(bg_color)
        except:
            pass

def add_navigation(slide, index, total, style_config, title=""):
    nav_config = style_config.get('layout', {}).get('navigation', False)
    if not nav_config:
        return
        
    # Basic navigation: "01. Introduction" top left
    left = Inches(0.5)
    top = Inches(0.2)
    width = Inches(8)
    height = Inches(0.5)
    
    # Creating a text box
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    
    p = tf.paragraphs[0]
    p.text = f"{index:02d}. {title.upper()}"
    p.font.size = Pt(10)
    
    accent_color = style_config.get('colors', {}).get('accent', '#000000')
    try:
        p.font.color.rgb = hex_to_rgb(accent_color)
    except:
        pass
        
    font_name = style_config.get('typography', {}).get('heading', {}).get('font', 'Arial')
    p.font.name = font_name

def parse_markdown(md_content):
    slides = []
    current_slide = None
    
    lines = md_content.split('\n')
    for line in lines:
        line = line.strip()
        if not line:
            continue
            
        # Headers -> New Slide
        if line.startswith('#'):
            # Detect level
            level = len(line.split()[0])
            title = line.lstrip('#').strip()
            
            if current_slide:
                slides.append(current_slide)
            
            # Default layout is 1 (Title + Content). Can be overridden.
            layout_idx = 1
            if len(slides) == 0:
                layout_idx = 0 # Title Slide
            elif level == 2: # ## Section Header
                layout_idx = 2 # Section Header

            current_slide = {
                'title': title, 
                'layout': layout_idx, 
                'content': [], 
                'notes': []
            }
        
        # Speaker Notes
        elif line.startswith('<!-- note:'):
             if current_slide:
                 note_text = line.replace('<!-- note:', '').replace('-->', '').strip()
                 current_slide['notes'].append(note_text)

        # Images -> Add to current slide
        elif line.startswith('![') and '](' in line:
            # ![alt](path)
            m = re.match(r'!\[(.*?)\]\((.*?)\)', line)
            if m and current_slide:
                image_path = m.group(2)
                current_slide['content'].append({'type': 'image', 'path': image_path})
        
        # Bullets -> Add to current slide
        elif line.startswith('- ') or line.startswith('* '):
            if current_slide:
                text = line[1:].strip()
                current_slide['content'].append({'type': 'bullet', 'text': text})
        
        # Text -> Add as paragraph or handle special tags
        else:
             if current_slide:
                if line.startswith('<!-- layout:'): # Manual layout override
                     # <!-- layout: 0 --> or <!-- layout: Title Only -->
                     layout_val = line.replace('<!-- layout:', '').replace('-->', '').strip()
                     if layout_val.isdigit():
                         current_slide['layout'] = int(layout_val)
                else:
                    current_slide['content'].append({'type': 'text', 'text': line})

    if current_slide:
        slides.append(current_slide)
            
    return slides

def create_ppt(slides, output_file, template_file=None, style_config=None):
    if template_file and os.path.exists(template_file):
        prs = Presentation(template_file)
    else:
        prs = Presentation()

    total_slides = len(slides)

    for i, slide_data in enumerate(slides):
        layout_idx = slide_data.get('layout', 1)
        
        if layout_idx >= len(prs.slide_layouts):
            layout_idx = 1 # Fallback
            
        slide_layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(slide_layout)

        # Apply Style to Background
        if style_config:
            apply_slide_background(slide, style_config)
            # Add Navigation (Skip Title Slide 0)
            if layout_idx != 0:
                add_navigation(slide, i+1, total_slides, style_config, slide_data['title'])

        # Title
        if slide.shapes.title:
            slide.shapes.title.text = slide_data['title']
            if style_config:
                apply_text_style(slide.shapes.title, style_config, 'title')

        # Speaker Notes
        if slide_data['notes']:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = "\n".join(slide_data['notes'])

        # Body/Content Shape Finding
        body_shape = None
        if layout_idx == 0: # Title Slide
            if len(slide.placeholders) > 1:
                body_shape = slide.placeholders[1] # Subtitle
        else:
             for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:
                    body_shape = shape
                    break
             if not body_shape and len(slide.placeholders) > 1:
                body_shape = slide.placeholders[1]

        # Populate Content
        if body_shape and body_shape.has_text_frame:
            tf = body_shape.text_frame
            tf.text = "" 
            
            first_p = True
            for item in slide_data['content']:
                if item['type'] in ['bullet', 'text']:
                    if first_p:
                        p = tf.paragraphs[0]
                        first_p = False
                    else:
                        p = tf.add_paragraph()
                    
                    p.text = item['text']
                    if item['type'] == 'bullet':
                         p.level = 0
                    
                    # Apply specific text styling overrides if possible, otherwise rely on shape-level apply_text_style
                
                elif item['type'] == 'image':
                    left = Inches(1)
                    top = Inches(2)
                    height = Inches(4)
                    try:
                        slide.shapes.add_picture(item['path'], left, top, height=height)
                    except Exception as e:
                        print(f"Error adding image {item['path']}: {e}")

            if style_config:
                apply_text_style(body_shape, style_config, 'body')

    prs.save(output_file)

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python md2pptx.py <input.md> <output.pptx> [--template <template.pptx>] [--config <config.yaml>]")
        sys.exit(1)
        
    input_md = sys.argv[1]
    output_pptx = sys.argv[2]
    
    template = None
    config_file = None
    
    # Simple argument parsing
    args = sys.argv[3:]
    if '--template' in args:
        idx = args.index('--template')
        if idx + 1 < len(args):
            template = args[idx+1]
            
    if '--config' in args:
        idx = args.index('--config')
        if idx + 1 < len(args):
            config_file = args[idx+1]

    style_config = {}
    if config_file and os.path.exists(config_file):
        with open(config_file, 'r', encoding='utf-8') as f:
            style_config = yaml.safe_load(f)
            print(f"Loaded style config from {config_file}")

    with open(input_md, 'r', encoding='utf-8') as f:
        md_content = f.read()

    slides = parse_markdown(md_content)
    create_ppt(slides, output_pptx, template, style_config)
    print(f"Generated {output_pptx} with {len(slides)} slides.")
